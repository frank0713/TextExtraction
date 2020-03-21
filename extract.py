# _*_ coding: UTF-8 _*_
# @version: python 3.7.4
# @author: frank0713


import sqlite3
import os
import time
# pip install python-docx
from docx import Document
# pip install python-pptx
from pptx import Presentation
# pip install pypiwin32
import win32com.client as wc
import pandas as pd
from PIL import Image
# pip install pytesseract
# pip install tesseract-ocr
import pytesseract
# pip install pdf2image
from pdf2image import convert_from_path
from pdfminer.pdfparser import PDFParser, PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LTTextBoxHorizontal, LAParams
# uncompress
import zipfile
import tarfile
# pip install py7zr
import py7zr


# ToDo: move to another file
# ToDo: create dir to store
def create_db(db_name):
    """create a sqlite database"""

    db_names = db_name + ".db"
    conn = sqlite3.connect(db_names)
    table_name = db_name
    c = conn.cursor()
    c.execute("CREATE TABLE %s (Name TEXT, Extension TEXT, CTime TEXT, \
    MTime TEXT, Path TEXT, Size TEXT, Text TEXT)" % table_name)
    conn.commit()
    conn.close()


def data2db(file_dir, db_name):
    """push the data into sqlite database"""

    information = file_information(file_dir)
    db_names = str(db_name) + '.db'
    table_name = db_name
    conn = sqlite3.connect(db_names)
    c = conn.cursor()
    for i in information:
        c.execute("INSERT INTO %s (Name, Extension, CTime, MTime, Path,\
         Size, Text) VALUES ('%s', '%s', '%s', '%s', '%s', '%s', '%s')"
                  % (table_name, i['Name'], i['Extension'], i['CTime'],
                     i['MTime'], i['Path'], i['Size'], i['Text']))
        conn.commit()
    conn.close()


def timestamp2time(timestamp):
    """convert time to structured format"""

    time_stamp = time.localtime(timestamp)
    time_str = time.strftime('%Y-%m-%d %H:%M:%S', time_stamp)

    return time_str


def word(file):
    """extract text from .doc .docx based on docx module"""

    doc = Document(file)
    paras = ''
    for p in doc.paragraphs:
        para = p.text
        paras = paras + ' ' + para

    return paras


# ToDo: 最后索引完，删除最后一个temp.docx
# ToDo: SaveAs路径修改
def doc2docx(file):
    """convert .doc to .docx"""

    wordapp = wc.Dispatch("Word.Application")
    doc = wordapp.Documents.Open(file)
    doc.SaveAs(r"d:\\textgps\\temp.docx", 12)
    doc.Close()
    wordapp.Quit()


# ToDo: 最后索引完，删除最后一个temp.pptx
# ToDo: SaveAs路径修改
# ToDo: 解决弹窗
def ppt2pptx(file):
    """convert .ppt to .pptx"""

    pptapp = wc.Dispatch("PowerPoint.Application")
    ppt = pptapp.Presentations.Open(file)
    ppt.SaveAs(r"d:\\textgps\\temp.pptx")
    ppt.Close()
    pptapp.Quit()


def txt(file):
    """extract text from .txt files excluding .csv"""

    with open(file, 'r', encoding='utf-8') as f:
        text = f.read()

    return text


# ToDo: 最后删除temp.txt文件
def excel(file):
    """extract text from .xls/.xlsx files"""

    df = pd.read_excel(file, header=None)
    df.to_csv('temp.txt', header=None, sep=' ', index=False)
    text = txt('temp.txt')

    return text


# ToDo: 最后删除temp.txt文件
def csv(file):
    """extract text from .csv files"""

    df = pd.read_csv(file, header=None)
    df.to_csv('temp.txt', header=None, sep=' ', index=False)
    text = txt('temp.txt')

    return text


def pptx(file):
    """extract text from .pptx files"""

    shape_ts = []
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                t = shape.text
                shape_ts.append(t)
    text = ' '.join(shape_ts)

    return text


# ToDo: 最后删除temp.jpg文件
# ToDo: ocr语言
def scanpdf2txt(file, direction):
    """extract text from scanned pdf files by OCR using tesseract"""

    text = []
    pages = convert_from_path(file, 500)
    image_counter = 1
    for p in pages:
        fn = "temp.jpg"
        p.save(fn, 'JPEG')
        image_counter += 1
        img_path = os.path.join(direction, fn)
        t = pytesseract.image_to_string(Image.open(img_path), lang='chi_sim')
        text.append(t)
    texts = ' '.join(text)

    return texts


def docpdf2txt(file):
    """extract text from doc pdf files"""
    
    with open(file, 'rb') as fp:
        praser = PDFParser(fp)
        pdf = PDFDocument()
        praser.set_document(pdf)
        pdf.set_parser(praser)
    rsrcmgr = PDFResourceManager()
    laparams = LAParams()
    device = PDFPageAggregator(rsrcmgr, laparams=laparams)
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    text = []
    for p in pdf.get_pages():
        interpreter.process_page(p)
        layout = device.get_result()
        for x in layout:
            if isinstance(x, LTTextBoxHorizontal):
                t = x.get_text()
                text.append(t)
    texts = ' '.join(text)

    return texts


# ToDo: make a dir for store temp_zip files
# ToDo: 最后删除temp_zip整个文件夹
def unzip(file):
    """uncompress .zip file temporally"""

    zip_file = zipfile.ZipFile(file)
    # os.mkdir('temp_zip')
    for f in zip_file.namelist():
        zip_file.extract(f, 'D:\\temp_zip\\')
    zip_file.close()


# ToDo: make a dir for store temp_zip files
# ToDo: 最后删除temp_zip整个文件夹
def untar(file):
    """uncompress .tar file temporally"""

    tar_file = tarfile.open(file)
    for f in tar_file.getnames():
        tar_file.extract(f, 'D:\\temp_zip\\')
    tar_file.close()


# ToDo: make a dir for store temp_zip files
# ToDo: 最后删除temp_zip整个文件夹
def un7z(file):
    """uncompress .7z file temporally"""

    sevenz_file = py7zr.SevenZipFile(file, mode='r')
    sevenz_file.extractall(path='D:\\temp_zip\\')
    sevenz_file.close()


def uncompress2txt():
    zip_info = file_information('D:\\temp_zip\\')  # zip_info = list
    
    return zip_info


# ToDo: Class(file_information)
def file_information(file_dir):
    """fetch info of customized format files under selected path: file_dir"""

    information = []
    # ToDo: change to an optional model
    word_new_format = ['.docx']
    word_old_format = ['.doc']  # .doc需转为.docx
    excel_format = ['.xlsx', '.xls']
    csv_format = ['.csv']
    txt_format = ['.txt']
    ppt_new_format = ['.pptx']
    ppt_old_format = ['.ppt']  # .ppt需转为.pptx
    pdf_format = ['.pdf']
    zip_format = ['.zip']
    sevenz_format = ['.7z']
    tar_format = ['.tar']

    for root, dirs, files in os.walk(file_dir):
        for f in files:
            path = os.path.join(root, f)
            extension = os.path.splitext(f)[1]
            name = os.path.splitext(f)[0]
            ctime = timestamp2time(os.path.getctime(path))
            mtime = timestamp2time(os.path.getmtime(path))
            size = os.path.getsize(path) / float(1024 * 1024)  # Unit: Mb
            if extension in word_new_format:
                paras = word(file=path)
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': paras}
                information.append(info)
            if extension in word_old_format:
                doc2docx(path)
                paras = word(r"d:\\textgps\\temp.docx")
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': paras}
                information.append(info)
            if extension in txt_format:
                text = txt(file=path)
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': text}
                information.append(info)
            if extension in excel_format:
                text = excel(file=path)
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': text}
                information.append(info)
            if extension in csv_format:
                text = csv(file=path)
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': text}
                information.append(info)
            if extension in ppt_new_format:
                text = pptx(file=path)
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': text}
                information.append(info)
            if extension in ppt_old_format:
                ppt2pptx(path)
                text = pptx(r"d:\\textgps\\temp.pptx")
                info = {'Name': name, 'Extension': extension, 'CTime': ctime,
                        'MTime': mtime, 'Path': path, 'Size': size,
                        'Text': text}
                information.append(info)
            if extension in pdf_format:
                text = docpdf2txt(path)
                # 判定是doc还是scan pdf
                if text != '':
                    info = {'Name': name, 'Extension': extension,
                            'CTime': ctime, 'MTime': mtime, 'Path': path,
                            'Size': size, 'Text': text}
                    information.append(info)
                # ToDo: direction设为某个指定的隐藏临时文件夹，同其他
                else:
                    text = scanpdf2txt(path, direction='D:\\textgps\\')
                    info = {'Name': name, 'Extension': extension,
                            'CTime': ctime, 'MTime': mtime, 'Path': path,
                            'Size': size, 'Text': text}
                    information.append(info)
            if extension in zip_format:
                unzip(path)
                zip_info = uncompress2txt()
                for i in zip_info:
                    i['Name'] = i['Name'] + i['Extension'] + ' (in) ' + name
                    i['Path'] = path
                    i['Extension'] = '.zip'
                    information.append(i)
            if extension in tar_format:
                untar(path)
                tar_info = uncompress2txt()
                for i in tar_info:
                    i['Name'] = i['Name'] + i['Extension'] + ' (in) ' + name
                    i['Path'] = path
                    i['Extension'] = '.tar'
                    information.append(i)
            if extension in sevenz_format:
                un7z(path)
                sevenz_info = uncompress2txt()
                for i in sevenz_info:
                    i['Name'] = i['Name'] + i['Extension'] + ' (in) ' + name
                    i['Path'] = path
                    i['Extension'] = '.7z'
                    information.append(i)

    return information


create_db('test')
data2db('d:\\textgps', 'test')
