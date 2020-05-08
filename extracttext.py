# _*_ coding: UTF-8 _*_
# @version: python 3.7.4
# @author: frank0713

# ToDo: encode & decode
# ToDo: handle exceptions
# ToDo: 修改初始化路径，暂用用户主文件夹

import os
from subprocess import call
from shutil import rmtree
from chardet import detect
# **MS Office**
from comtypes.client import CreateObject  # .doc
# from win32com.client import DispatchEx
from docx import Document  # .docx
from pptx import Presentation  # .pptx
from pandas import read_csv, read_excel, read_table  # .csv/excel
# **MS Office exceptions**
from _ctypes import COMError  # .doc/.ppt
from docx.opc import exceptions  # .docx
from pptx import exc  # .pptx
from xlrd.biffh import XLRDError  # excel
from pandas.errors import EmptyDataError  # .csv
# **scanned pdf**
# from PIL import Image
from pdf2image import convert_from_path
from pdf2image.exceptions import PDFPageCountError
# import pytesseract
# **doc pdf**
# use xpdf in command line mode
# from pdfplumber import open as pdfopen
# from pdfminer.pdfdocument import PDFPasswordIncorrect
# **Markup formats**
from xml.etree.ElementTree import ElementTree, ParseError  # xml
from bs4 import BeautifulSoup  # html/htm
from tex2py import tex2py  # tex
# chm: use 7z uncompress to htm files, and extract
# **ODF**
from odf.opendocument import load as odfload
from odf import text as odftext
from odf import teletype
# **email**
from email.header import decode_header, make_header
from email import message_from_file  # eml
from extract_msg import Message  # msg
# **ebook**
from ebooklib import epub, ITEM_DOCUMENT
# mobi/azw/azw3: use kindleunpack.py to unpack, and extract


class TXTText:
    """
    Extract texts from general .txt file.
    
    It can also be used as a txt extractor for some other files, such as: 
    Markup language files(Markdown/Yaml); Code script files(C++/Python...etc.).

    MS-Excel associated files(.xls/.xlsx/.xlsm/.csv) convert into .txt format, 
    using Pandas library, to be extracted.
    """

    def __init__(self, path):
        self.path = path

    def txttext(self):
        """Extract texts from general .txt files."""

        try:
            with open(self.path, 'rb') as f:
                unicode_text = f.read()
                code = detect(unicode_text)['encoding']  # detect code type
                text = unicode_text.decode(encoding=code)
                text = text.replace("'", "‘")
        except TypeError:  # decode error or empty file(cannot detect code type)
            text = ''
        return text

    @staticmethod
    def ctxttext(file):
        """Extract texts from converted temporal .txt file for other inner class
        methods."""

        try:
            with open(file, 'rb') as f:
                unicode_text = f.read()
                code = detect(unicode_text)['encoding']  # detect code type
                print(code)
                text = unicode_text.decode(encoding=code)
                text = text.replace("'", "‘")
        except TypeError:  # decode error or empty file(cannot detect code type)
            text = ''
        return text

    @staticmethod
    def initialize_txt_path():
        """Initialize path for temporal converted .txt files."""

        user_main_path = os.path.expanduser('~')
        directory = user_main_path + '\\Appdata\\Local\\Temp\\Textgps\\temp_txt'
        if not os.path.exists(directory):
            os.makedirs(directory)
        return directory

    def csvtext(self):
        """Extract texts from .csv files. Use Pandas library to convert into .txt format."""

        try:
            df = read_csv(self.path, header=None)
            # ToDo: 修改创建dir。暂用用户主文件夹下路径
            directory = self.initialize_txt_path()
            csv2txt_file = directory + '\\temp.txt'
            df.to_csv(csv2txt_file, header=None, sep=' ', index=False)
            text = self.ctxttext(file=csv2txt_file)
        except (EmptyDataError, UnicodeDecodeError):
            # empty .csv(also: locked file) or decode error
            text = ''
        return text

    def tsvtext(self):
        """Extract texts from .tsv files. Use Pandas library to convert into .txt format."""

        try:
            df = read_table(self.path, header=None)
            # ToDo: 修改创建dir。暂使用用户主文件夹下路径
            directory = self.initialize_txt_path()
            tsv2txt_file = directory + '\\temp.txt'
            df.to_csv(tsv2txt_file, header=None, sep='\t', index=False)
            text = self.ctxttext(file=tsv2txt_file)
        except (EmptyDataError, UnicodeDecodeError):
            # empty .tsv(also: locked file) or decode error
            text = ''
        return text

    def exceltext(self):
        """Extract texts from .xls/.xlsx/.xlsm files.

        Use Pandas library to convert the MS-Excel associated format files
        (including old version[.xls] or macro[.xlsm] ones) into .txt format."""

        try:
            df = read_excel(self.path, header=None)
            # ToDo: 修改创建路径。暂使用用户主文件夹下路径
            directory = self.initialize_txt_path()
            excel2txt_file = directory + '\\.temp.txt'
            df.to_csv(excel2txt_file, header=None, sep=' ', index=False)
            text = self.ctxttext(file=excel2txt_file)
        except XLRDError:  # empty or locked file
            text = ''
        return text

    @staticmethod
    def rm_txt_files():
        """Remove the temporal .txt files at the end."""

        txt_file = TXTText.initialize_txt_path() + 'temp.txt'
        if os.path.exists(txt_file):
            os.remove(txt_file)


class WordText:
    """
    Extract texts from .docx/.doc/.docm/.rtf files.

    For .docx files, use docx library.

    For old version MS-Word files(.doc), macro-Word files(.docm), or rich text
    format files(.rtf), use comtypes.client module.
    """

    def __init__(self, path):
        self.path = path

    def docxtext(self):
        """Extract texts from .docx files."""

        try:
            doc = Document(self.path)
        except exceptions.PackageNotFoundError:
            # cannot open file: locked file, empty file...
            text = ''
        else:
            # main body text
            text = ''
            for p in doc.paragraphs:
                para = p.text
                text = text + para + ' '
            # table text
            table_text = []
            for table in doc.tables:
                for row in range(0, len(table.rows)):
                    r_t = []
                    for column in range(0, len(table.columns)):
                        t = table.cell(row, column).text
                        r_t.append(t)
                    r_t = ' '.join(r_t)
                    table_text.append(r_t)
            table_text = ' \n'.join(table_text)
            text = text + table_text
            text = text.replace("'", "‘")
        return text

    def doctext(self):
        """Extract texts from .doc/.docm/.rtf files, using comtypes.client."""

        try:
            wordapp = CreateObject("Word.Application")
            doc = wordapp.Documents.Open(self.path, PasswordDocument=' ')
        except COMError:  # locked file
            text = ''
        else:
            texts = []
            for para in doc.paragraphs:
                t = para.Range.Text
                texts.append(t)
            text = ' '.join(texts)
            text = text.replace("'", "‘")
            doc.Close()
            wordapp.Quit()
        return text


class PPTText:
    """
    Extract texts from .pptx/.ppt/.pptm files.

    For .pptx files, use pptx library.

    For old version MS-PowerPoint files(.ppt) or macro-PowerPoint files(.pptm), 
    use comtypes.client module.
    """

    def __init__(self, path):
        self.path = path

    def pptxtext(self):
        """Extract texts from .pptx files."""

        try:
            shape_ts = []
            ppt = Presentation(self.path)
        except exc.PackageNotFoundError:
            # cannot open file: locked or empty file
            shape_ts = []
        else:
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        t = shape.text
                        shape_ts.append(t)
        text = ' '.join(shape_ts)
        text = text.replace("'", "‘")
        return text

    def ppttext(self):
        """Extract texts from .ppt/.pptm files, using comtypes.client module."""

        try:
            pptapp = CreateObject("PowerPoint.Application")
            pwd = ' '
            path_pwd = str(self.path) + "::" + pwd
            ppt = pptapp.Presentations.Open(path_pwd, WithWindow=False)
        except COMError:  # locked file
            text = ''
        else:
            texts = []
            slide_count = ppt.Slides.Count
            for i in range(1, slide_count + 1):
                shape_count = ppt.Slides(i).Shapes.Count
                for j in range(1, shape_count + 1):
                    if ppt.Slides(i).Shapes(j).HasTextFrame:
                        t = ppt.Slides(i).Shapes(j).TextFrame.TextRange.Text
                        texts.append(t)
            text = ' '.join(texts)
            text = text.replace("'", "‘")
            ppt.Close()
            pptapp.Quit()
        return text


class ODFText:
    """Extract texts from ODF files(.odt/.ods/.odp)."""

    def __init__(self, path):
        self.path = path

    def odftext(self):
        """Extract texts from .odt/.ods/.odp files"""

        odf_file = odfload(self.path)
        odf_text = odf_file.getElementsByType(odftext.P)
        text = ''
        for para in odf_text:
            t = teletype.extractText(para)
            text = text + t + ' '
        text = text.replace("'", "‘")
        return text


class MarkupText:
    """
    Extract texts from markup format files(.xml/.html/.tex/.chm)

    For .md(markdown) or .yml(yaml) files, use TXT extractor(class TXTText) directly.
    """

    def __init__(self, path):
        self.path = path

    def xmltext(self):
        """Extract texts from .xml files."""

        with open(self.path, 'rb') as dxml:
            unicode_text = dxml.read()
            code = detect(unicode_text)['encoding']  # detect code type
            # print(code)
            if code == "None":  # empty file or cannot detect code typy
                text = ''
            else:
                try:
                    with open(file=self.path, encoding=code) as xf:
                        tree = ElementTree(file=xf)
                        root = tree.getroot()
                        texts = []
                        for child in root.iter():
                            t = child.text
                            texts.append(t)
                        text = ' '.join(texts)
                        text = text.replace("'", "‘")
                except ParseError:  # wrong code or others...
                    text = ''
        return text

    def htmltext(self):
        """Extract texts from .html files."""

        with open(self.path, 'rb') as dhf:
            unicode_text = dhf.read()
            code = detect(unicode_text)['encoding']
            # print(code)
            if code == "None":  # empty file or connot detect code type
                text = ''
            else:
                try:
                    with open(self.path, 'r', encoding=code) as hf:
                        html = BeautifulSoup(hf, "html.parser")
                        text = html.get_text()
                        text = text.replace("'", "‘")
                except AttributeError:  # wrong code or others...
                    text = ''
        return text

    def textext(self):
        """Extract texts from .tex files."""

        with open(self.path, 'rb') as dtf:
            unicode_text = dtf.read()
            code = detect(unicode_text)['encoding']
            print(code)
            if code == "None":  # empty file or cannot detect code type
                text = ''
            else:
                try:
                    with open(self.path, 'r', encoding=code) as tf:
                        toc = tex2py(tf.read())
                        text = []
                        for i in toc.descendants:
                            if isinstance(i, str):
                                text.append(i)
                        text = ' '.join(text)
                        text = text.replace("'", "‘")
                except UnicodeDecodeError:  # wrong code type
                    text = ''
        return text

    @staticmethod
    def initial_chm_path():
        """
        Initialize the path for temporal unzipped chm files.
        """

        user_main_path = os.path.expanduser('~')
        chm_dir = user_main_path + '\\Appdata\\Local\\Temp\\Textgps\\temp_chm'
        if not os.path.exists(chm_dir):
            os.makedirs(chm_dir)
        return chm_dir

    def chmtext(self):
        """
        Uncompress the .chm file, and extract the texts from output .htm files.
        """

        # ToDo: 修改7zip路径
        # ToDo: 修改初始化路径，暂用用户主文件夹
        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        output_dir = os.path.join(self.initial_chm_path(), file_name)
        output_path = '-o' + output_dir
        sevenzip_path = "D:\\textgps\\7z.exe"
        cmd = [sevenzip_path, 'x', self.path, '*.htm', '-y', output_path, '-r']
        # 7zip path, extract mode, input file, uncompressed format, all yes,
        # output path, recursively uncompress
        call(cmd)
        # extract texts from htm files
        texts = []
        for root, dirs, files in os.walk(output_dir):
            for f in files:
                f_path = os.path.join(root, f)
                with open(f_path, 'rb') as dhf:
                    unicode_text = dhf.read()
                    code = detect(unicode_text)['encoding']
                    # print(code)
                    if code == "None":  # empty file or connot detect code type
                        t = ''
                        texts.append(t)
                    else:
                        try:
                            with open(f_path, 'r', encoding=code) as hf:
                                html = BeautifulSoup(hf, "html.parser")
                                t = html.get_text()
                                t = t.replace("'", "‘")
                                texts.append(t)
                        except AttributeError:  # wrong code or others...
                            t = ''
                            texts.append(t)
        text = ' '.join(texts)
        return text

    def rm_chm(self):
        """
        Remove the temporal uncompressed chm files at the end.
        """

        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        output_dir = os.path.join(self.initial_chm_path(), file_name)
        if os.path.exists(output_dir):
            rmtree(output_dir)


class PDFText:
    """
    Extract texts from .pdf files.

    One type is document style, using xpdf to convert to .txt; the other is
    scanned type, which is converted to image and extracted by OCR(tesseract).
    """

    def __init__(self, path):
        self.path = path

    @staticmethod
    def initialize_dpdf_path():
        """Initialize path for temporal doc-pdf files."""

        # ToDo: 修改初始化路径，暂用用户主文件夹
        user_main_path = os.path.expanduser('~')
        dpdf_dir = user_main_path + '\\Appdata\\Local\\Temp\\Textgps\\temp_dpdf'
        if not os.path.exists(dpdf_dir):
            os.makedirs(dpdf_dir)
        return dpdf_dir

    @staticmethod
    def initialize_spdf_path():
        """Initialize path for temporal scan-pdf files."""

        # ToDo: 修改初始化路径，暂用用户主文件夹
        user_main_path = os.path.expanduser('~')
        spdf_dir = user_main_path + '\\Appdata\\Local\\Temp\\Textgps\\temp_spdf'
        if not os.path.exists(spdf_dir):
            os.makedirs(spdf_dir)
        return spdf_dir

    def docpdftext(self):
        """Extract texts from document type PDF."""

        # ToDo: 修改xpdf路径
        xpdf_path = 'D:\\xpdf\\pdftotext.exe'
        directory = self.initialize_dpdf_path()
        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        convert_path = directory + file_name + '.txt'
        call([xpdf_path, '-enc', 'UTF-8', self.path, convert_path])
        try:
            with open(convert_path, 'r', encoding='utf-8') as f:
                text = f.read()
                text = text.replace("'", "‘")
        except FileNotFoundError:  # locked file cannot be converted, so it is not found
            text = ''
        return text

    def scanpdftext(self):
        """Extract texts from scanned PDF by OCR(tesseract)."""

        # ToDo: 修改初始化路径，暂用用户主文件夹
        directory = self.initialize_spdf_path()
        texts = []
        try:
            pages = convert_from_path(self.path, 300)
            for p in pages:
                fn = os.path.join(directory, "temp.jpg")
                p.save(fn, "JPEG")
                cmd = ['d:\\programs\\tesseract\\tesseract.exe',
                       fn, directory, '--tessdata-dir',
                       'd:\\programs\\tesseract\\tessdata', '-l', 'chi_sim+eng',
                       '--dpi', '300', '--oem', '1']
                # app path, image path, output dir, tessdata path, language(in sequence),
                # dpi value setting, OCR engine mode(1: Automatic page segmentation with OSD)
                # Input "APP-PATH --help-extra" in cmd to get the details.
                call(cmd)
                out_txt = directory + '\\temp.txt'
                with open(out_txt, 'r', encoding='utf-8') as f:
                    t = f.read()
                texts.append(t)
            text = ' '.join(texts)
            text = text.replace("'", "‘")
        except (ValueError, PDFPageCountError):  # cannot parser or others...
            text = ''
        return text

    @staticmethod
    def rm_pdf():
        """Remove the temporal files of scanned pdf extractor at the end."""

        sdpf_path = PDFText.initialize_spdf_path()
        dpdf_path = PDFText.initialize_dpdf_path()
        if os.path.exists(sdpf_path):
            for f in os.listdir(sdpf_path):
                os.remove(f)
        if os.path.exists(dpdf_path):
            for f in os.listdir(sdpf_path):
                os.remove(f)


class MailText:
    """
    Extract texts and some infomation from email files.
    """

    def __init__(self, path):
        self.path = path

    def emltext(self):
        """
        Extract texts and some information from .eml format email files, using builtin module: email.
        """

        texts = []
        with open(self.path, 'r') as eml:
            message = message_from_file(eml)
        send = 'From: ' + message['From']
        to = 'To: ' + message['To']
        dt = 'DateTime: ' + message['Date']
        try:
            sub = make_header(decode_header(message['Subject']))
            header = 'Subject: ' + str(sub)
        except UnicodeDecodeError:  # wrong decode
            header = ''
        texts.append(send)
        texts.append(to)
        texts.append(dt)
        texts.append(header)
        # mail content
        texts.append('Message: ')
        try:
            for part in message.walk():
                if not part.is_multipart():
                    text_type = part.get_content_subtype()
                    msgs = part.get_payload(decode=True)
                    code_type = part.get_content_charset()
                    # 避免繁体中文解码错误
                    if code_type == 'gb2312':
                        code_type = 'gbk'
                    if text_type == 'plain':
                        try:
                            tmail_text = msgs.decode(code_type)
                        except (UnicodeDecodeError, AttributeError):
                            # wrong decode or empty text
                            tmail_text = ''
                        texts.append(tmail_text)
                    elif text_type == 'html':
                        try:
                            hmail = BeautifulSoup(msgs.decode(code_type),
                                                  "html.parser")
                            hmail_text = hmail.body.get_text()
                        except (UnicodeDecodeError, AttributeError):
                            # wrong decode or empty text
                            hmail_text = ''
                        texts.append(hmail_text)
        except (UnicodeDecodeError, AttributeError):  # wrong decode or empty text
            msgs = ''
            texts.append(msgs)
        text = ', '.join(texts)
        text = text.replace("'", "‘")
        return text

    def msgtext(self):
        """
        Extract texts and some information from .msg format email files.
        """

        try:
            mail = Message(self.path)
        except (AttributeError, Exception):  # wrong decode or cannot parse...
            text = ''
        else:
            if mail.sender is None:
                send = 'From: '
            else:
                send = 'From: ' + mail.sender
            if mail.to is None:
                to = 'To: '
            else:
                to = 'To: ' + mail.to
            if mail.date is None:
                dt = 'DataTime: '
            else:
                dt = 'DateTime: ' + mail.date
            if mail.subject is None:
                sub = 'Subject: '
            else:
                sub = 'Subject: ' + mail.subject
            if mail.body is None:
                msgs = 'Message: '
            else:
                msgs = 'Message: ' + mail.body
            text = send + ', ' + to + ', ' + dt + ', ' + sub + ', ' + msgs
            text = text.replace("'", "‘")
        return text


class EbookText:
    """
    Extract texts from .epub/.mobi/.azw/.azw3 files.
    """

    def __init__(self, path):
        self.path = path

    def epubtext(self):
        """
        Extract texts from .epub files.
        """

        with open(self.path, 'rb') as f:
            ebook = epub.read_epub(f)
        texts = []
        for doc in ebook.get_items_of_type(ITEM_DOCUMENT):
            try:
                body = doc.get_body_content().decode('utf-8')
                body_parser = BeautifulSoup(body, 'html.parser')
                try:
                    t = body_parser.get_text()
                except AttributeError:  # empty text
                    t = ''
                texts.append(t)
            except UnicodeDecodeError:  # wrong decode
                t = ''
                texts.append(t)
        text = ' '.join(texts)
        text = text.replace("'", "‘")
        return text

    @staticmethod
    def initial_mobi_path():
        """
        Initialize the path for temporal mobi files unpacked by kindleunpack.py
        """

        # ToDo: 修改初始化路径，暂用用户主文件夹
        user_main_path = os.path.expanduser('~')
        mobi_dir = user_main_path + '\\Appdata\\Local\\Temp\\Textgps\\temp_mobi'
        if not os.path.exists(mobi_dir):
            os.makedirs(mobi_dir)
        return mobi_dir

    @staticmethod
    def initial_azw3_path():
        """
        Initialize the path for temporal azw files unpacked by kindleunpack.py
        """

        # ToDo: 修改初始化路径，暂用用户主文件夹
        user_main_path = os.path.expanduser('~')
        azw_dir = user_main_path + '\\Appdata\\Local\\Temp\\Textgps\\temp_azw3'
        if not os.path.exists(azw_dir):
            os.makedirs(azw_dir)
        return azw_dir

    def mobitext(self):
        """
        Extract texts from .mobi/azw files, using kindleunpack.py to uncompress the
        mobi/azw file and extract texts from a html file.

        Note: azw is different from azw3, azw and mobi formats are also called mobi7 or
        earlier version mobi file. However, azw is commonly protected by DRM.
        This function just extracts the texts from unencrypted azw file(without DRM),
        it would not crack DRM.
        """

        # ToDo: 修改kindleunpack.py的路径
        # ToDo: 修改临时文件路径，暂用用户主文件夹
        directory = self.initial_mobi_path()
        kindelunpack_path = 'd:\\ku\\lib\\kindleunpack.py'
        in_dir = self.path
        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        out_dir = os.path.join(directory, file_name)
        # ToDo: python cmd支持
        cmd = ['python', kindelunpack_path, '-d', in_dir, out_dir]
        # python app, kindleunpack.py path, dump the file, input file, output file
        call(cmd)
        html_path = os.path.join(out_dir, 'mobi7/book.html')
        try:
            text = MarkupText(path=html_path).htmltext()
        except FileNotFoundError:  # DRM file or cannot unpack properly...
            text = ''
        return text

    def azw3text(self):
        """
        Extract texts from azw3 file, using kindleunpack.py to uncompress the azw3 file
        and extract texts from an epub file.

        Note: azw3 format is also called mobi8, commonly it is protected by DRM. This function
        just extracts the texts from unencrypted azw3 file(without DRM) and it would not
        crack DRM.
        """

        # ToDo: 修改kindleunpack.py的路径
        # ToDo: 修改临时文件路径，暂用用户主文件夹
        directory = self.initial_azw3_path()
        kindelunpack_path = 'd:\\ku\\lib\\kindleunpack.py'
        in_dir = self.path
        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        out_dir = os.path.join(directory, file_name)
        # ToDo: python cmd支持
        cmd = ['python', kindelunpack_path, '-d', in_dir, out_dir]
        # python app, kindleunpack.py path, dump the file, input file, output file
        call(cmd)
        epub_path = out_dir + '/mobi8/' + file_name + '.epub'
        # extract texts from the same name epub file in epub_path
        try:
            with open(epub_path, 'rb') as f:
                ebook = epub.read_epub(f)
            texts = []
            for doc in ebook.get_items_of_type(ITEM_DOCUMENT):
                try:
                    body = doc.get_body_content().decode('utf-8')
                    body_parser = BeautifulSoup(body, 'html.parser')
                    try:
                        t = body_parser.get_text()
                    except AttributeError:  # empty text
                        t = ''
                    texts.append(t)
                except UnicodeDecodeError:  # wrong decode
                    t = ''
                    texts.append(t)
            text = ' '.join(texts)
            text = text.replace("'", "‘")
        except FileNotFoundError:  # DRM file or unpack inproperly...
            text = ''
        return text

    def rm_mobi(self):
        """
        Remove the temporal mobi unpacked files at the end.
        """

        out_dir = self.initial_mobi_path()
        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        mobi_path = os.path.join(out_dir, file_name)
        if os.path.exists(mobi_path):
            rmtree(mobi_path)

    def rm_azw3(self):
        """
        Remove the temporal azw3 unpacked files at the end.
        """

        out_dir = self.initial_azw3_path()
        file_name = os.path.splitext(os.path.split(self.path)[-1])[0]
        azw3_path = os.path.join(out_dir, file_name)
        if os.path.exists(azw3_path):
            rmtree(azw3_path)
